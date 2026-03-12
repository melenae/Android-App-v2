from pptx import Presentation
from pptx.util import Pt


def add_title_slide(prs: Presentation, title: str, subtitle: str | None = None) -> None:
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    if subtitle:
        subtitle_placeholder.text = subtitle


def add_bullet_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
) -> None:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame

    title_placeholder.text = title

    if not bullets:
        return

    # Первый пункт
    tf.text = bullets[0]
    tf.paragraphs[0].font.size = Pt(20)

    # Остальные пункты
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)


def build_presentation() -> Presentation:
    prs = Presentation()

    # Титульный слайд
    add_title_slide(
        prs,
        title="ERP‑Tools",
        subtitle="Операционная система для управления развитием корпоративных IT‑систем",
    )

    # Слайд 1 — проблема рынка
    add_bullet_slide(
        prs,
        title="Почему ERP‑проекты проваливаются",
        bullets=[
            "Требования хранятся в Excel и почте",
            "Документация устаревает сразу после внедрения",
            "Критические знания находятся в головах сотрудников",
            "Любое изменение может сломать систему",
            "Задачи не связаны с архитектурой и бизнес‑процессами",
            "80% времени IT‑команда тратит на координацию, а не на результат",
        ],
    )

    # Слайд 2 — реальная картина в компаниях
    add_bullet_slide(
        prs,
        title="Реальная картина: зоопарк инструментов",
        bullets=[
            "Jira — задачи",
            "Excel — требования и реестр доработок",
            "Word — ТЗ и спецификации",
            "Wiki / Confluence — документация",
            "Service Desk — обращения пользователей",
            "Отдельные диаграммы и схемы архитектуры",
            "Инструменты не связаны между собой",
        ],
    )

    # Слайд 3 — главная проблема
    add_bullet_slide(
        prs,
        title="Главная проблема",
        bullets=[
            "Нет единой модели IT‑системы компании",
            "Изменения опасны и непредсказуемы",
            "Знания теряются при уходе ключевых людей",
            "Внедрение и развитие идут медленно",
            "Новые сотрудники долго входят в контекст",
        ],
    )

    # Слайд 4 — решение
    add_bullet_slide(
        prs,
        title="Решение: ERP‑Tools",
        bullets=[
            "Единая среда управления IT‑развитием компании",
            "В одной системе: архитектура, требования, задачи, roadmap, support, документация",
            "Связь от бизнес‑процесса до внедрённого изменения",
            "Single source of truth для IT‑ландшафта компании",
        ],
    )

    # Слайд 5 — ключевая идея
    add_bullet_slide(
        prs,
        title="Ключевая идея",
        bullets=[
            "Связка сквозного цикла развития:",
            "Бизнес‑процесс → Функция",
            "→ Требование → Задача",
            "→ Разработка → Тест → Внедрение",
            "Прозрачная трассировка от идеи до результата",
        ],
    )

    # Слайд 6 — 3 модуля
    add_bullet_slide(
        prs,
        title="3 ключевых модуля ERP‑Tools",
        bullets=[
            "1️⃣ Управление roadmap IT‑развития",
            "2️⃣ Функциональное моделирование и архитектура",
            "3️⃣ Service Desk и поддержка пользователей",
        ],
    )

    # Слайд 7 — ключевые преимущества
    add_bullet_slide(
        prs,
        title="Ключевые преимущества",
        bullets=[
            "Вся архитектура IT описана в единой модели",
            "Все изменения контролируемы и отслеживаемы",
            "Знания о системе не теряются",
            "Внедрение и модернизация идут быстрее",
            "Уменьшается риск провала ERP‑проектов",
        ],
    )

    # Слайд 8 — цифры (примерные)
    add_bullet_slide(
        prs,
        title="Измеримый эффект (примерные цифры)",
        bullets=[
            "−50% времени на модернизацию ERP и смежных систем",
            "+40% прозрачность процессов и статуса изменений",
            "+25% производительность IT‑команды",
            "Снижение количества инцидентов из‑за ошибок изменений",
        ],
    )

    # Слайд 9 — кому это нужно
    add_bullet_slide(
        prs,
        title="Кому особенно нужен ERP‑Tools",
        bullets=[
            "Компании с ERP‑системами (1С, SAP, Oracle и др.)",
            "IT‑департаменты средних и крупных компаний",
            "ERP‑интеграторы и консалтинговые компании",
            "Офисы цифровой трансформации",
        ],
    )

    # Слайд 10 — пример внедрения
    add_bullet_slide(
        prs,
        title="Пример внедрения: до и после",
        bullets=[
            "До ERP‑Tools:",
            "Excel / Word — разрозненные требования и ТЗ",
            "Jira / Trello — задачи без связи с архитектурой",
            "Wiki — неполная и устаревающая документация",
            "",
            "После внедрения ERP‑Tools:",
            "Единая модель IT‑ландшафта",
            "Связанные между собой архитектура, требования, задачи и сервис",
            "Прозрачный контроль изменений и знаний",
        ],
    )

    return prs


def main() -> None:
    prs = build_presentation()
    output_path = "ERP-Tools_presentation.pptx"
    prs.save(output_path)
    print(f"Презентация успешно сохранена в файле: {output_path}")


if __name__ == "__main__":
    main()

